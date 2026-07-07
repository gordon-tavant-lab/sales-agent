#!/usr/bin/env python3
"""
pptx_screenshot.py — Convert PPTX slides to PNG images.

Usage:
    python3 pptx_screenshot.py <file.pptx> [output_dir] [slide_range]

Examples:
    python3 pptx_screenshot.py deck.pptx
    python3 pptx_screenshot.py deck.pptx ./out/
    python3 pptx_screenshot.py deck.pptx ./out/ 1-5
    python3 pptx_screenshot.py deck.pptx ./out/ 1,3,7

Renderer priority:
    1. LibreOffice headless (install: brew install --cask libreoffice)
    2. python-pptx + Pillow fallback renderer
"""

import os
import sys
import shutil
import subprocess
import tempfile
from pathlib import Path


# ── Helpers ──────────────────────────────────────────────────────────────────

def parse_slide_range(spec: str, total: int) -> list[int]:
    """Parse '1-5' or '1,3,7' into 0-based slide indices."""
    if not spec:
        return list(range(total))
    indices = set()
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            indices.update(range(int(a) - 1, int(b)))
        else:
            indices.add(int(part) - 1)
    return sorted(i for i in indices if 0 <= i < total)


def find_libreoffice():
    for cmd in ("soffice", "libreoffice",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if shutil.which(cmd) or Path(cmd).exists():
            return cmd
    return None


# ── Renderer 1: LibreOffice ───────────────────────────────────────────────────

def render_libreoffice(pptx_path: Path, output_dir: Path, slide_indices: list[int]) -> list[Path]:
    """Use LibreOffice to convert PPTX → PNG (best quality)."""
    soffice = find_libreoffice()
    with tempfile.TemporaryDirectory() as tmp:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "png", "--outdir", tmp, str(pptx_path)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice failed: {result.stderr}")

        # LibreOffice names files: <stem>.png (only first page) or <stem>1.png, <stem>2.png...
        tmp_path = Path(tmp)
        stem = pptx_path.stem
        pngs = sorted(tmp_path.glob(f"{stem}*.png"))

        if not pngs:
            raise RuntimeError("LibreOffice produced no PNG files")

        out_files = []
        for idx, png in enumerate(pngs):
            if slide_indices and idx not in slide_indices:
                continue
            dest = output_dir / f"slide-{idx + 1:03d}.png"
            shutil.copy2(png, dest)
            out_files.append(dest)

    return out_files


# ── Renderer 2: python-pptx + Pillow fallback ─────────────────────────────────

def rgb_from_color(color_obj):
    """Extract RGB tuple from a pptx color, handling theme colors."""
    try:
        if color_obj and color_obj.rgb:
            h = str(color_obj.rgb)
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        pass
    return None


def emu_to_px(emu: int, scale: float) -> int:
    return int(emu / 9144 * scale)


def render_fallback(pptx_path: Path, output_dir: Path, slide_indices: list[int]) -> list[Path]:
    """Render slides using python-pptx + Pillow. Shows colors, text, and shapes."""
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN
    from PIL import Image, ImageDraw, ImageFont
    import textwrap

    RENDER_W = 1280
    RENDER_H = 720

    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    scale_x = RENDER_W / slide_w
    scale_y = RENDER_H / slide_h

    # Try to load a system font
    def get_font(size: int, bold: bool = False):
        font_paths = [
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
            "/System/Library/Fonts/SFNSDisplay.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]
        for fp in font_paths:
            if Path(fp).exists():
                try:
                    return ImageFont.truetype(fp, size)
                except Exception:
                    pass
        return ImageFont.load_default()

    out_files = []
    for idx, slide in enumerate(prs.slides):
        if slide_indices and idx not in slide_indices:
            continue

        # Background color
        bg_color = (255, 255, 255)
        try:
            bg_fill = slide.background.fill
            if bg_fill.type is not None:
                c = rgb_from_color(bg_fill.fore_color)
                if c:
                    bg_color = c
        except Exception:
            pass

        img = Image.new("RGB", (RENDER_W, RENDER_H), bg_color)
        draw = ImageDraw.Draw(img)

        # Slide number watermark (small, top-right)
        draw.text((RENDER_W - 60, 10), f"#{idx + 1}", fill=(180, 180, 180), font=get_font(14))

        # Draw shapes
        try:
            from pptx.util import Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            for shape in slide.shapes:
                try:
                    # Position and size
                    x = int(shape.left * scale_x) if shape.left else 0
                    y = int(shape.top * scale_y) if shape.top else 0
                    w = int(shape.width * scale_x) if shape.width else 0
                    h = int(shape.height * scale_y) if shape.height else 0

                    # Draw shape background fill
                    try:
                        fill = shape.fill
                        if fill.type is not None:
                            fc = rgb_from_color(fill.fore_color)
                            if fc and w > 0 and h > 0:
                                draw.rectangle([x, y, x + w, y + h], fill=fc)
                    except Exception:
                        pass

                    # Draw shape border
                    try:
                        line = shape.line
                        if line.color and line.color.type is not None:
                            lc = rgb_from_color(line.color)
                            if lc and w > 0 and h > 0:
                                draw.rectangle([x, y, x + w, y + h], outline=lc, width=2)
                    except Exception:
                        pass

                    # Draw text
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                text = run.text.strip()
                                if not text:
                                    continue
                                # Font size
                                try:
                                    font_size = int(run.font.size / 12700 * scale_y * 0.75) if run.font.size else 12
                                    font_size = max(8, min(font_size, 48))
                                except Exception:
                                    font_size = 12
                                # Text color
                                tc = (0, 0, 0)
                                try:
                                    tc = rgb_from_color(run.font.color) or (0, 0, 0)
                                except Exception:
                                    pass
                                font = get_font(font_size)
                                try:
                                    draw.text((x + 4, y + 4), text[:80], fill=tc, font=font)
                                    y += font_size + 2
                                except Exception:
                                    pass
                except Exception:
                    continue
        except Exception:
            pass

        dest = output_dir / f"slide-{idx + 1:03d}.png"
        img.save(dest, "PNG")
        out_files.append(dest)
        print(f"  ✓ slide-{idx + 1:03d}.png")

    return out_files


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    pptx_path = Path(sys.argv[1]).expanduser().resolve()
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    output_dir = Path(sys.argv[2]).expanduser().resolve() if len(sys.argv) > 2 else pptx_path.parent / "screenshots"
    slide_range_spec = sys.argv[3] if len(sys.argv) > 3 else ""

    output_dir.mkdir(parents=True, exist_ok=True)

    # Count slides to parse range
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    slide_indices = parse_slide_range(slide_range_spec, total)

    print(f"📊 {pptx_path.name} — {total} slides total, exporting {len(slide_indices)}")
    print(f"📁 Output: {output_dir}")

    soffice = find_libreoffice()
    if soffice:
        print(f"🖥  Renderer: LibreOffice ({soffice})")
        try:
            files = render_libreoffice(pptx_path, output_dir, slide_indices)
        except Exception as e:
            print(f"⚠️  LibreOffice failed ({e}), falling back to python-pptx renderer")
            files = render_fallback(pptx_path, output_dir, slide_indices)
    else:
        print("🐍 Renderer: python-pptx + Pillow (install LibreOffice for pixel-perfect output)")
        print("   brew install --cask libreoffice")
        files = render_fallback(pptx_path, output_dir, slide_indices)

    print(f"\n✅ Done — {len(files)} PNG files written to {output_dir}")
    for f in files[:5]:
        print(f"   {f.name}")
    if len(files) > 5:
        print(f"   ... and {len(files) - 5} more")


if __name__ == "__main__":
    main()
